import streamlit as st
import fitz                          # PyMuPDF  — PDF annotations
from docx import Document            # python-docx — DOCX read/write
from docx.shared import Pt, RGBColor
from pptx import Presentation        # python-pptx — PPTX notes/comments
from ebooklib import epub, ITEM_DOCUMENT  # ebooklib — EPUB <mark> tags
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm
from reportlab.lib import colors
from bs4 import BeautifulSoup        # parse EPUB HTML
import re, os
from io import BytesIO

# ═══════════════════════════════════════════════
#  PAGE CONFIG
# ═══════════════════════════════════════════════

st.set_page_config(
    page_title="Highlight Extractor",
    page_icon="📑",
    layout="centered",
    initial_sidebar_state="expanded"
)

st.markdown(
    '<link href="https://fonts.googleapis.com/css2?family=DM+Serif+Display:ital@0;1'
    '&family=DM+Sans:wght@300;400;500;600&display=swap" rel="stylesheet">',
    unsafe_allow_html=True
)

# ═══════════════════════════════════════════════
#  CSS
# ═══════════════════════════════════════════════

st.markdown("""
<style>
:root {
    --ac: #E8572A; --ac-l: #FF7A52;
    --ac-bg: rgba(232,87,42,0.10); --ac-br: rgba(232,87,42,0.30);
    --r-md: 12px; --r-lg: 18px;
}
html, body, [class*="css"] { font-family: 'DM Sans', sans-serif !important; }
#MainMenu, footer, header { visibility: hidden; }
.stDeployButton { display: none; }
.block-container { padding: 2rem 1.75rem 4rem !important; max-width: 700px !important; }

/* HERO */
.hero { text-align: center; padding: 2rem 0 0.5rem; }
.badge {
    display: inline-block; background: var(--ac-bg); border: 1px solid var(--ac-br);
    color: var(--ac); font-size: 10px; font-weight: 700; letter-spacing: 0.13em;
    text-transform: uppercase; padding: 5px 16px; border-radius: 999px; margin-bottom: 1.1rem;
}
.h-title {
    font-family: 'DM Serif Display', Georgia, serif !important;
    font-size: 2.8rem; font-weight: 400; line-height: 1.15;
    margin: 0 0 0.6rem; letter-spacing: -0.02em;
}
.h-title em { font-style: italic; color: var(--ac); }
.h-sub { font-size: 14px; font-weight: 300; opacity: 0.55; margin: 0 auto; max-width: 420px; line-height: 1.7; }
.div-line { height: 1px; background: rgba(128,128,128,0.15); margin: 1.5rem 0 1.75rem; }

/* SECTION LABEL */
.sec-label {
    font-size: 11px; font-weight: 700; letter-spacing: 0.1em;
    text-transform: uppercase; opacity: 0.4; margin-bottom: 4px;
}

/* UPLOAD ZONE */
[data-testid="stFileUploaderDropzone"] {
    border: 1.5px dashed var(--ac-br) !important; border-radius: var(--r-lg) !important;
    background: var(--ac-bg) !important; padding: 2rem 1.5rem !important; transition: border-color 0.2s;
}
[data-testid="stFileUploaderDropzone"]:hover { border-color: var(--ac) !important; }
[data-testid="stFileUploaderDropzone"] button {
    background: var(--ac) !important; color: white !important; border: none !important;
    border-radius: 8px !important; font-size: 13px !important; font-weight: 600 !important;
    padding: 0.45rem 1.2rem !important;
}
[data-testid="stFileUploaderDropzone"] small { font-size: 12px !important; opacity: 0.45 !important; }

/* FORMAT PICKER */
.fmt-grid {
    display: grid; grid-template-columns: repeat(4, 1fr);
    gap: 8px; margin: 0.5rem 0 1.25rem;
}
.fmt-card {
    border: 1.5px solid rgba(128,128,128,0.15); border-radius: var(--r-md);
    padding: 0.75rem 0.5rem; text-align: center; cursor: pointer;
    transition: all 0.15s; user-select: none;
}
.fmt-card:hover { border-color: var(--ac-br); background: var(--ac-bg); }
.fmt-card.active { border-color: var(--ac); background: var(--ac-bg); }
.fmt-icon { font-size: 20px; margin-bottom: 4px; }
.fmt-name { font-size: 11px; font-weight: 600; letter-spacing: 0.05em; }
.fmt-desc { font-size: 10px; opacity: 0.45; margin-top: 2px; }

/* CARDS */
.icard {
    border-radius: var(--r-md); border: 1px solid rgba(128,128,128,0.13);
    padding: 1rem 1.2rem; margin: 1.1rem 0; display: flex; align-items: center; gap: 12px;
}
.icard-icon { font-size: 22px; flex-shrink: 0; }
.icard-name { font-weight: 600; font-size: 14px; margin: 0; }
.icard-meta { font-size: 12px; opacity: 0.45; margin: 3px 0 0; }
.icard-status { margin-left: auto; font-size: 10px; font-weight: 700; letter-spacing: 0.1em; text-transform: uppercase; opacity: 0.35; }

/* STAT GRID */
.stat-grid { display: grid; grid-template-columns: 1fr 1fr; gap: 10px; margin: 1.2rem 0; }
.stat-card { border-radius: var(--r-md); border: 1px solid rgba(128,128,128,0.11); padding: 1rem 1.2rem; text-align: center; }
.stat-n { font-family: 'DM Serif Display', Georgia, serif !important; font-size: 2.4rem; color: var(--ac); line-height: 1; display: block; margin-bottom: 6px; }
.stat-l { font-size: 11px; font-weight: 600; letter-spacing: 0.08em; text-transform: uppercase; opacity: 0.4; }

/* SUCCESS BANNER */
.ok-banner {
    border-radius: var(--r-md); border: 1px solid rgba(34,197,94,0.28);
    background: rgba(34,197,94,0.08); padding: 0.9rem 1.2rem;
    display: flex; align-items: center; gap: 12px; margin: 0.75rem 0 1rem;
}
.ok-dot { width: 32px; height: 32px; border-radius: 50%; background: rgba(34,197,94,0.18); display: flex; align-items: center; justify-content: center; font-size: 16px; flex-shrink: 0; }
.ok-title { font-weight: 600; font-size: 14px; color: #16a34a; margin: 0 0 2px; }
.ok-sub { font-size: 12px; opacity: 0.6; margin: 0; }

/* DOWNLOAD BUTTON */
.stDownloadButton button {
    width: 100% !important; background: var(--ac) !important; color: white !important;
    border: none !important; border-radius: var(--r-md) !important;
    padding: 0.8rem 1.5rem !important; font-size: 15px !important; font-weight: 600 !important;
    box-shadow: 0 3px 16px rgba(232,87,42,0.28) !important; transition: background 0.18s, transform 0.12s !important;
}
.stDownloadButton button:hover { background: var(--ac-l) !important; transform: translateY(-1px) !important; }

/* RADIO — style as pill tabs */
div[data-testid="stRadio"] > div { display: flex; gap: 8px; flex-wrap: wrap; }
div[data-testid="stRadio"] label {
    border: 1.5px solid rgba(128,128,128,0.18) !important; border-radius: 999px !important;
    padding: 6px 16px !important; font-size: 13px !important; cursor: pointer;
    transition: all 0.15s !important;
}
div[data-testid="stRadio"] label:has(input:checked) {
    border-color: var(--ac) !important; background: var(--ac-bg) !important; color: var(--ac) !important;
}

/* SIDEBAR */
[data-testid="stSidebar"] { border-right: 1px solid rgba(128,128,128,0.1) !important; }
[data-testid="stSidebar"] .block-container { padding: 1.75rem 1.1rem !important; }
.s-logo { display: flex; align-items: center; gap: 10px; margin-bottom: 1.75rem; }
.s-icon { width: 34px; height: 34px; background: var(--ac); border-radius: 9px; display: flex; align-items: center; justify-content: center; font-size: 17px; }
.s-name { font-family: 'DM Serif Display', Georgia, serif !important; font-size: 15px; line-height: 1.25; }
.s-name small { font-family: 'DM Sans', sans-serif !important; font-size: 11px; opacity: 0.4; font-weight: 400; display: block; }
.s-sec { font-size: 10px; font-weight: 700; letter-spacing: 0.1em; text-transform: uppercase; opacity: 0.35; margin: 1.4rem 0 0.6rem; }
.s-row { display: flex; gap: 8px; align-items: flex-start; margin-bottom: 0.75rem; font-size: 13px; line-height: 1.5; }
.s-num { flex-shrink: 0; width: 21px; height: 21px; background: var(--ac-bg); border: 1px solid var(--ac-br); border-radius: 50%; display: flex; align-items: center; justify-content: center; font-size: 10px; font-weight: 700; color: var(--ac); margin-top: 1px; }
.s-chip { display: inline-block; background: rgba(128,128,128,0.1); border-radius: 4px; font-size: 10px; font-weight: 700; padding: 2px 7px; letter-spacing: 0.05em; margin: 2px 2px 0 0; }
.s-foot { font-size: 11.5px; opacity: 0.4; line-height: 1.65; margin-top: 1.5rem; padding-top: 1.2rem; border-top: 1px solid rgba(128,128,128,0.12); }
.empty { text-align: center; padding: 1rem 0 0.5rem; font-size: 13px; opacity: 0.35; }
</style>
""", unsafe_allow_html=True)

# ═══════════════════════════════════════════════
#  SIDEBAR
# ═══════════════════════════════════════════════

with st.sidebar:
    st.markdown("""
    <div class="s-logo">
        <div class="s-icon">📑</div>
        <div class="s-name">Highlight Extractor<small>Multi-format → Any output</small></div>
    </div>

    <div class="s-sec">Supported inputs</div>
    <div class="s-row"><div><span class="s-chip">PDF</span> Highlight annotations</div></div>
    <div class="s-row"><div><span class="s-chip">DOCX</span> Highlighted / commented runs</div></div>
    <div class="s-row"><div><span class="s-chip">PPTX</span> Slide notes &amp; comments</div></div>
    <div class="s-row"><div><span class="s-chip">EPUB</span> &lt;mark&gt; tagged spans</div></div>
    <div class="s-row"><div><span class="s-chip">TXT</span> Full text extraction</div></div>

    <div class="s-sec">Output formats</div>
    <div class="s-row"><div><span class="s-chip">DOCX</span> Structured Word document</div></div>
    <div class="s-row"><div><span class="s-chip">PDF</span> Portable document</div></div>
    <div class="s-row"><div><span class="s-chip">TXT</span> Plain text, one line each</div></div>
    <div class="s-row"><div><span class="s-chip">MD</span> Markdown with headings</div></div>

    <div class="s-sec">How it works</div>
    <div class="s-row"><div class="s-num">1</div><div><strong>Upload</strong> any supported file</div></div>
    <div class="s-row"><div class="s-num">2</div><div><strong>Pick</strong> your output format</div></div>
    <div class="s-row"><div class="s-num">3</div><div><strong>Extract</strong> — duplicates auto-removed</div></div>
    <div class="s-row"><div class="s-num">4</div><div><strong>Download</strong> your clean notes</div></div>

    <div class="s-foot">ALL-CAPS highlights automatically become section headings in the output.</div>
    """, unsafe_allow_html=True)

# ═══════════════════════════════════════════════
#  HERO
# ═══════════════════════════════════════════════

st.markdown("""
<div class="hero">
    <div class="badge">✦ Multi-format Annotation Tool</div>
    <div class="h-title">Extract your <em>highlights</em>,<br>your way.</div>
    <div class="h-sub">Upload PDF, DOCX, PPTX, EPUB, or TXT — get clean notes
    in whatever format you need.</div>
</div>
<div class="div-line"></div>
""", unsafe_allow_html=True)

# ═══════════════════════════════════════════════
#  FILE UPLOAD
# ═══════════════════════════════════════════════

st.markdown('<div class="sec-label">Upload your file</div>', unsafe_allow_html=True)

uploaded_file = st.file_uploader(
    label="upload",
    type=["pdf", "docx", "pptx", "epub", "txt"],
    label_visibility="collapsed",
    help="Supports PDF, DOCX, PPTX, EPUB, TXT. Max 200MB."
)

# ═══════════════════════════════════════════════
#  HELPERS — TEXT CLEANING
# ═══════════════════════════════════════════════

def clean_text(text):
    text = text.replace("\n", " ")
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

def add_unique(seen, items, new_text):
    t = clean_text(new_text)
    if t and t not in seen:
        seen.add(t)
        items.append(t)

# ═══════════════════════════════════════════════
#  EXTRACTORS — one per input type
# ═══════════════════════════════════════════════

def extract_pdf(file_bytes, filename):

    path = filename

    with open(path, "wb") as f:
        f.write(file_bytes)

    pdf = fitz.open(path)

    seen = set()
    items = []
    pages = 0

    for page in pdf:

        page_had = False

        annotations = page.annots()

        if annotations:

            for annot in annotations:

                # Only highlights
                if annot.type[1] == "Highlight":

                    try:

                        quad_points = annot.vertices

                        words = page.get_text("words")

                        highlighted_words = []

                        # Loop through every quad
                        for i in range(0, len(quad_points), 4):

                            quad = quad_points[i:i+4]

                            rect = fitz.Quad(quad).rect

                            for word in words:

                                word_rect = fitz.Rect(word[:4])

                                # Check intersection
                                if rect.intersects(word_rect):

                                    highlighted_words.append(word[4])

                        highlighted_text = " ".join(highlighted_words)

                        highlighted_text = clean_text(
                            highlighted_text
                        )

                        if (
                            highlighted_text
                            and highlighted_text not in seen
                        ):

                            seen.add(highlighted_text)

                            items.append(highlighted_text)

                            page_had = True

                    except Exception:
                        pass

        if page_had:
            pages += 1

    return items, pages

def extract_docx(file_bytes, filename):
    path = filename
    with open(path, "wb") as f:
        f.write(file_bytes)
    doc = Document(path)
    seen, items = set(), []
    # highlighted runs (shading / highlight colour set)
    for para in doc.paragraphs:
        for run in para.runs:
            rpr = run._r.get_or_add_rPr()
            highlight = rpr.find(
                "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}highlight"
            )
            if highlight is not None:
                add_unique(seen, items, run.text)
    # comments via XML
    try:
        comments_part = doc.part.package.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments"
        )
        from lxml import etree
        ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        for comment in comments_part._element.findall(f"{{{ns}}}comment"):
            text = "".join(t.text or "" for t in comment.iter(f"{{{ns}}}t"))
            add_unique(seen, items, text)
    except Exception:
        pass
    return items, len(items)

def extract_pptx(file_bytes, filename):
    path = filename
    with open(path, "wb") as f:
        f.write(file_bytes)
    prs = Presentation(path)
    seen, items = set(), []
    slides_with_content = 0
    for slide in prs.slides:
        had = False
        # Slide notes
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                add_unique(seen, items, notes)
                had = True
        # Comments
        try:
            for comment in slide.comments:
                add_unique(seen, items, comment.text)
                had = True
        except Exception:
            pass
        if had:
            slides_with_content += 1
    return items, slides_with_content

def extract_epub(file_bytes, filename):
    path = filename
    with open(path, "wb") as f:
        f.write(file_bytes)
    book = epub.read_epub(path)
    seen, items = set(), []
    for item in book.get_items_of_type(ITEM_DOCUMENT):
        soup = BeautifulSoup(item.get_content(), "html.parser")
        for tag in soup.find_all(["mark", "span"]):
            cls = " ".join(tag.get("class", []))
            style = tag.get("style", "")
            if tag.name == "mark" or "highlight" in cls.lower() or "background" in style.lower():
                add_unique(seen, items, tag.get_text())
    return items, len(items)

def extract_txt(file_bytes, filename):
    text = file_bytes.decode("utf-8", errors="ignore")
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    seen, items = set(), []
    for line in lines:
        add_unique(seen, items, line)
    return items, len(items)

# ═══════════════════════════════════════════════
#  BUILDERS — one per output type
# ═══════════════════════════════════════════════

def build_docx(items, title):
    doc = Document()
    doc.add_heading(title, level=1)
    for item in items:
        if item.isupper():
            doc.add_heading(item, level=2)
        else:
            doc.add_paragraph(item)
    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document", ".docx"

def build_pdf(items, title):

    from reportlab.platypus import (
        SimpleDocTemplate,
        Paragraph,
        Spacer,
        PageBreak,
        KeepTogether
    )

    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus.tables import Table, TableStyle
    from reportlab.platypus.flowables import HRFlowable

    buf = BytesIO()

    doc = SimpleDocTemplate(
        buf,
        pagesize=A4,
        leftMargin=1.8*cm,
        rightMargin=1.8*cm,
        topMargin=1.8*cm,
        bottomMargin=1.8*cm,
    )

    styles = getSampleStyleSheet()

    title_style = ParagraphStyle(
        "TitleStyle",
        parent=styles["Title"],
        fontName="Helvetica-Bold",
        fontSize=22,
        leading=28,
        textColor=colors.HexColor("#E8572A"),
        spaceAfter=20,
    )

    heading_style = ParagraphStyle(
        "HeadingStyle",
        parent=styles["Heading2"],
        fontName="Helvetica-Bold",
        fontSize=14,
        leading=20,
        textColor=colors.HexColor("#222222"),
        spaceBefore=14,
        spaceAfter=8,
    )

    body_style = ParagraphStyle(
        "BodyStyle",
        parent=styles["BodyText"],
        fontName="Helvetica",
        fontSize=11,
        leading=18,
        textColor=colors.HexColor("#333333"),
        spaceAfter=10,
    )

    story = []

    # Title
    story.append(Paragraph(title, title_style))
    story.append(HRFlowable(width="100%"))
    story.append(Spacer(1, 0.3 * cm))

    for item in items:

        # Clean dangerous HTML chars
        safe = (
            item.replace("&", "&amp;")
                .replace("<", "&lt;")
                .replace(">", "&gt;")
        )

        # Remove weird unicode bullets/emojis if needed
        safe = safe.encode("utf-8", "ignore").decode("utf-8")

        # Heading detection
        if item.isupper():

            story.append(Spacer(1, 0.15 * cm))
            story.append(Paragraph(safe, heading_style))

        else:

            bullet_text = f"• {safe}"

            story.append(
                Paragraph(
                    bullet_text,
                    body_style
                )
            )

    doc.build(story)

    pdf_bytes = buf.getvalue()

    buf.close()

    return (
        pdf_bytes,
        "application/pdf",
        ".pdf"
    )

def build_txt(items, title):
    lines = [title, "=" * len(title), ""]
    for item in items:
        lines.append(item)
        lines.append("")
    content = "\n".join(lines)
    return content.encode("utf-8"), "text/plain", ".txt"

def build_md(items, title):
    lines = [f"# {title}", ""]
    for item in items:
        if item.isupper():
            lines.append(f"## {item}")
        else:
            lines.append(f"- {item}")
        lines.append("")
    content = "\n".join(lines)
    return content.encode("utf-8"), "text/markdown", ".md"

# ═══════════════════════════════════════════════
#  FILE TYPE METADATA
# ═══════════════════════════════════════════════

FILE_META = {
    "pdf":  ("📕", "PDF Document",       "Highlight annotations"),
    "docx": ("📘", "Word Document",      "Highlighted runs & comments"),
    "pptx": ("📙", "PowerPoint File",    "Slide notes & comments"),
    "epub": ("📗", "eBook (EPUB)",       "<mark> tagged spans"),
    "txt":  ("📄", "Plain Text File",    "All lines extracted"),
}

OUTPUT_OPTIONS = {
    "DOCX — Word Document":  ("docx", "📘"),
    "PDF — Portable Document": ("pdf", "📕"),
    "TXT — Plain Text":      ("txt", "📄"),
    "MD — Markdown":         ("md",  "📝"),
}

# ═══════════════════════════════════════════════
#  MAIN FLOW
# ═══════════════════════════════════════════════

if uploaded_file is not None:

    ext = uploaded_file.name.rsplit(".", 1)[-1].lower()
    size_kb = round(uploaded_file.size / 1024, 1)
    size_str = f"{size_kb} KB" if size_kb < 1024 else f"{round(size_kb/1024,1)} MB"
    icon, ftype, fhint = FILE_META.get(ext, ("📄", "File", "Text extraction"))

    st.markdown(f"""
    <div class="icard">
        <div class="icard-icon">{icon}</div>
        <div>
            <div class="icard-name">{uploaded_file.name}</div>
            <div class="icard-meta">{size_str} · {ftype} · {fhint}</div>
        </div>
        <div class="icard-status">Ready</div>
    </div>
    """, unsafe_allow_html=True)

    # ── Output format picker ──────────────────────
    st.markdown('<div class="sec-label" style="margin-top:1rem">Choose output format</div>', unsafe_allow_html=True)

    out_choice = st.radio(
        label="output_format",
        options=list(OUTPUT_OPTIONS.keys()),
        index=0,
        horizontal=True,
        label_visibility="collapsed"
    )

    out_fmt, out_icon = OUTPUT_OPTIONS[out_choice]

    # ── Extract button ────────────────────────────
    st.markdown('<div style="margin-top:1rem"></div>', unsafe_allow_html=True)

    if st.button("⚡  Extract highlights", use_container_width=True):

        file_bytes = uploaded_file.read()
        base_name = os.path.splitext(uploaded_file.name)[0]

        with st.spinner("Extracting content…"):
            if ext == "pdf":
                items, pages = extract_pdf(file_bytes, uploaded_file.name)
            elif ext == "docx":
                items, pages = extract_docx(file_bytes, uploaded_file.name)
            elif ext == "pptx":
                items, pages = extract_pptx(file_bytes, uploaded_file.name)
            elif ext == "epub":
                items, pages = extract_epub(file_bytes, uploaded_file.name)
            elif ext == "txt":
                items, pages = extract_txt(file_bytes, uploaded_file.name)
            else:
                items, pages = [], 0

        total = len(items)
        page_label = "Slides" if ext == "pptx" else ("Lines" if ext == "txt" else "Pages")

        st.markdown(f"""
        <div class="stat-grid">
            <div class="stat-card">
                <span class="stat-n">{total}</span>
                <span class="stat-l">Items extracted</span>
            </div>
            <div class="stat-card">
                <span class="stat-n">{pages}</span>
                <span class="stat-l">{page_label} with content</span>
            </div>
        </div>
        """, unsafe_allow_html=True)

        if total > 0:
            title = f"{base_name} — Extracted Highlights"

            if out_fmt == "docx":
                data, mime, suffix = build_docx(items, title)
            elif out_fmt == "pdf":
                data, mime, suffix = build_pdf(items, title)
            elif out_fmt == "txt":
                data, mime, suffix = build_txt(items, title)
            else:
                data, mime, suffix = build_md(items, title)

            out_filename = f"{base_name}_Highlights{suffix}"

            st.markdown(f"""
            <div class="ok-banner">
                <div class="ok-dot">✓</div>
                <div>
                    <div class="ok-title">Ready to download</div>
                    <div class="ok-sub">{total} items → {out_icon} {out_choice.split('—')[0].strip()}</div>
                </div>
            </div>
            """, unsafe_allow_html=True)

            st.download_button(
                label=f"⬇  Download  {out_filename}",
                data=data,
                file_name=out_filename,
                mime=mime,
            )

        else:
            st.warning(
                f"⚠️ No extractable content found in this {ftype}. "
                "For PDFs, ensure highlights are saved as annotations. "
                "For DOCX, use Word's text highlight feature. "
                "For EPUB, look for <mark> tags in the source."
            )

else:
    st.markdown(
        '<div class="empty">Supports PDF · DOCX · PPTX · EPUB · TXT — drag above or click to browse</div>',
        unsafe_allow_html=True
    )